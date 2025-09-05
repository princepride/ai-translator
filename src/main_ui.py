import zipfile
from typing import Optional, Tuple
from docx.opc.oxml import parse_xml
import yaml
import os
import shutil
import gradio as gr
from gradio.utils import NamedString
from utils.path import get_models
from utils.cuda import get_gpu_info
import json
import time
import importlib.util
from modules.file import FileReaderFactory, ExcelFileWriter
from docx.opc.pkgreader import _SerializedRelationships, _SerializedRelationship
from docx import Document
import markdown
from bs4 import BeautifulSoup
from docx.oxml.ns import qn
from docx.table import Table
from docx.text.paragraph import Paragraph
from docx.shared import Inches
from pptx import Presentation
import re
from transformers import AutoTokenizer
import openpyxl
from openai import OpenAI
import pandas as pd
import zhconv
from concurrent.futures import ThreadPoolExecutor, as_completed
from dotenv import load_dotenv
import logging
import logging.handlers
from datetime import datetime
import traceback
import functools
load_dotenv()
# 获取当前脚本所在目录的绝对路径
script_dir = os.path.dirname(os.path.abspath(__file__))

# ==================== LOGGING SETUP ====================
def setup_logging():
    """
    Setup comprehensive logging system with hierarchical directory structure.
    Creates separate loggers for different categories and organizes logs by date.
    """
    # Create logs directory structure
    logs_base_dir = os.path.join(script_dir, '..', 'logs')
    current_date = datetime.now()
    year_dir = os.path.join(logs_base_dir, str(current_date.year))
    month_dir = os.path.join(year_dir, f"{current_date.month:02d}")

    # Create directories if they don't exist
    os.makedirs(month_dir, exist_ok=True)

    # Define log file paths
    date_str = current_date.strftime("%d")
    info_log_path = os.path.join(month_dir, f"{date_str}_info.log")
    error_log_path = os.path.join(month_dir, f"{date_str}_error.log")

    # Configure log format
    log_format = logging.Formatter(
        '%(asctime)s - %(name)s - %(levelname)s - %(funcName)s:%(lineno)d - %(message)s',
        datefmt='%Y-%m-%d %H:%M:%S'
    )

    # Remove existing handlers to avoid duplicates
    for handler in logging.root.handlers[:]:
        logging.root.removeHandler(handler)

    # Setup root logger
    logging.basicConfig(level=logging.DEBUG, handlers=[])

    # Create and configure loggers
    loggers = {}

    # 1. General application logger (INFO and DEBUG)
    app_logger = logging.getLogger('translator.app')
    app_logger.setLevel(logging.DEBUG)

    # Info file handler with rotation
    info_handler = logging.handlers.RotatingFileHandler(
        info_log_path, maxBytes=10*1024*1024, backupCount=5, encoding='utf-8'
    )
    info_handler.setLevel(logging.DEBUG)
    info_handler.setFormatter(log_format)
    app_logger.addHandler(info_handler)

    # Console handler for development
    console_handler = logging.StreamHandler()
    console_handler.setLevel(logging.INFO)
    console_handler.setFormatter(log_format)
    app_logger.addHandler(console_handler)

    loggers['app'] = app_logger

    # 2. Error logger (ERROR and CRITICAL only)
    error_logger = logging.getLogger('translator.error')
    error_logger.setLevel(logging.ERROR)

    error_handler = logging.handlers.RotatingFileHandler(
        error_log_path, maxBytes=10*1024*1024, backupCount=5, encoding='utf-8'
    )
    error_handler.setLevel(logging.ERROR)
    error_handler.setFormatter(log_format)
    error_logger.addHandler(error_handler)

    loggers['error'] = error_logger

    # 3. Translation process logger
    translation_logger = logging.getLogger('translator.translation')
    translation_logger.setLevel(logging.DEBUG)
    translation_logger.addHandler(info_handler)
    translation_logger.addHandler(console_handler)

    loggers['translation'] = translation_logger

    return loggers

# Initialize logging system
loggers = setup_logging()
app_logger = loggers['app']
error_logger = loggers['error']
translation_logger = loggers['translation']

def log_function_call(logger_name='app'):
    """
    Decorator to log function entry and exit with execution time.
    """
    def decorator(func):
        @functools.wraps(func)
        def wrapper(*args, **kwargs):
            logger = loggers.get(logger_name, app_logger)
            func_name = func.__name__

            # Log function entry
            logger.info(f"ENTER {func_name} - Args: {len(args)} args, {len(kwargs)} kwargs")

            start_time = time.time()
            try:
                result = func(*args, **kwargs)
                execution_time = time.time() - start_time
                logger.info(f"EXIT {func_name} - Success - Execution time: {execution_time:.2f}s")
                return result
            except Exception as e:
                execution_time = time.time() - start_time
                error_logger.error(f"ERROR in {func_name} - Execution time: {execution_time:.2f}s - Error: {str(e)}")
                error_logger.error(f"Traceback: {traceback.format_exc()}")
                raise
        return wrapper
    return decorator

# Log application startup
app_logger.info("="*50)
app_logger.info("AI Translator Application Starting")
app_logger.info(f"Script directory: {script_dir}")
app_logger.info("="*50)

# 构建baseConfig.yml和modelExplains.yml的绝对路径
file_path = os.path.join(script_dir, 'configs', 'baseConfig.yml')
tokenizer = AutoTokenizer.from_pretrained(os.path.join(script_dir, 'tokenzier'))

with open(file_path, 'r') as file:
    yaml_data = yaml.load(file, Loader=yaml.FullLoader)

available_gpus = get_gpu_info()
api_models = get_models(os.path.join(script_dir, 'models/API'))
local_models = get_models(os.path.join(script_dir, 'models/local'))
available_models = {**api_models, **local_models}
# available_languages = [] # 注释掉或者移除这行，因为语言选项由 update_choices 动态生成

default_model_name = "gpt-4o-mini" # 定义默认模型名称
default_original_language = "Chinese"  # 定义默认原始语言
default_target_language_single = "English" # 定义默认目标语言 (单选)
default_target_language_multi = ["English"] # 定义默认目标语言 (多选)

def update_row_selection(selected_value):
    if selected_value == "所有行":
        return gr.update(visible=False)
    else:
        return gr.update(visible=True)

# --- translate_excel, translate, translate_excel_folder, word_to_markdown, markdown_to_word, translate_markdown_folder, glossary_check 函数保持不变 ---
# (省略这些函数的代码以保持简洁)
@log_function_call('translation')
def translate_excel(input_file, start_row, end_row, start_column, target_column,
                    selected_model,
                    selected_lora_model, selected_gpu, batch_size, original_language, target_languages):
    start_time = time.time()
    file_path = input_file.name

    translation_logger.info(f"Starting Excel translation - File: {os.path.basename(file_path)}")
    translation_logger.info(f"Parameters - Rows: {start_row}-{end_row}, Target column: {target_column}, "
                           f"Result column: {start_column}, Model: {selected_model}, Batch size: {batch_size}")
    translation_logger.info(f"Languages - From: {original_language}, To: {target_languages}")

    try:
        reader, fp = FileReaderFactory.create_reader(file_path)
        translation_logger.info(f"Successfully created file reader for {os.path.basename(file_path)}")

        inputs = reader.extract_text(file_path, target_column, start_row, end_row)
        translation_logger.info(f"Extracted {len(inputs)} text entries for translation")

        outputs = translate(inputs, selected_model, selected_lora_model, selected_gpu, batch_size, original_language,
                            target_languages)
        translation_logger.info(f"Translation completed - Generated {len(outputs)} outputs")

        excel_writer = ExcelFileWriter()
        print("Finally processed number: ", len(outputs))
        output_file = excel_writer.write_text(file_path, outputs, start_column, start_row, end_row)
        translation_logger.info(f"Results written to file: {os.path.basename(output_file)}")

        end_time = time.time()
        total_time = int(end_time - start_time)
        translation_logger.info(f"Excel translation completed successfully in {total_time}s")

        return f"Total process time: {total_time}s", output_file

    except Exception as e:
        error_logger.error(f"Failed to translate Excel file {os.path.basename(file_path)}: {str(e)}")
        error_logger.error(f"Traceback: {traceback.format_exc()}")
        raise

@log_function_call('translation')
def translate(inputs, selected_model, selected_lora_model, selected_gpu, batch_size, original_language, target_languages):
    if isinstance(inputs, str):
        inputs = [inputs]

    translation_logger.info(f"Starting translation - Model: {selected_model}, GPU: {selected_gpu}")
    translation_logger.info(f"Input count: {len(inputs)}, Batch size: {batch_size}")
    translation_logger.info(f"Languages - From: {original_language}, To: {target_languages}")

    model_path = available_models.get(selected_model)
    if not model_path:
        error_msg = f"Model '{selected_model}' not found in available models."
        error_logger.error(error_msg)
        print(error_msg)
        return []

    translation_logger.info(f"Using model path: {model_path}")

    model_file_path = os.path.join(model_path, 'model.py')
    if not os.path.exists(model_file_path):
        error_msg = f"No model.py found in {model_path}"
        error_logger.error(error_msg)
        print(error_msg)
        return []

    translation_logger.info(f"Loading model from: {model_file_path}")

    spec = importlib.util.spec_from_file_location("model", model_file_path)
    if spec is None or spec.loader is None:
        error_msg = f"Could not load spec for model.py in {model_path}"
        error_logger.error(error_msg)
        print(error_msg)
        return []

    model_module = importlib.util.module_from_spec(spec)
    try:
        spec.loader.exec_module(model_module)
        translation_logger.info("Model module loaded successfully")
    except Exception as e:
        error_msg = f"Error executing module {model_file_path}: {e}"
        error_logger.error(error_msg)
        error_logger.error(f"Traceback: {traceback.format_exc()}")
        print(error_msg)
        return []

    outputs = []
    if hasattr(model_module, 'Model'):
        try:
            translation_logger.info("Initializing model instance")
            model = model_module.Model(model_path, selected_lora_model, selected_gpu)

            if hasattr(model, 'generate'):
                translation_logger.info("Starting model generation")
                outputs = model.generate(inputs, original_language, target_languages, batch_size)
                translation_logger.info(f"Model generation completed - Generated {len(outputs)} outputs")
            else:
                error_msg = "Model class does not have a 'generate' method."
                error_logger.error(error_msg)
                print(error_msg)
        except Exception as e:
            error_msg = f"Error instantiating or running model from {model_path}: {e}"
            error_logger.error(error_msg)
            error_logger.error(f"Traceback: {traceback.format_exc()}")
            print(error_msg)
    else:
        error_msg = "No Model class found in model.py."
        error_logger.error(error_msg)
        print(error_msg)

    return outputs

@log_function_call('translation')
def translate_excel_folder(input_folder, start_row, end_row, start_column, target_column, selected_model,
                            selected_lora_model, selected_gpu, batch_size, original_language, target_languages,
                            row_selection):
    start_time = time.time()

    if not input_folder:
        error_logger.error("No files uploaded for folder translation")
        return "No files uploaded", []

    folder_path = os.path.dirname(input_folder[0].name)
    translation_logger.info(f"Starting Excel folder translation - {len(input_folder)} files")
    translation_logger.info(f"Folder path: {folder_path}")
    translation_logger.info(f"Parameters - Rows: {start_row}-{end_row}, Model: {selected_model}, Row selection: {row_selection}")

    processed_files = []
    processed_folder = os.path.join(folder_path, 'processed')
    os.makedirs(processed_folder, exist_ok=True)
    translation_logger.info(f"Created processed folder: {processed_folder}")

    for input_file in input_folder:
        file_path = input_file.name
        file_name = os.path.basename(file_path)
        translation_logger.info(f"Processing file: {file_name}")

        try:
            reader, updated_file_path = FileReaderFactory.create_reader(file_path)
            translation_logger.info(f"Successfully created reader for {file_name}")
        except ValueError as e:
            error_msg = f"Error creating reader for {file_name}: {e}"
            error_logger.error(error_msg)
            print(error_msg)
            continue
        original_file_obj_name = input_file.name # Store original name if needed
        if file_path != updated_file_path:
                file_path_to_process = updated_file_path
                temp_file_obj = NamedString(name=updated_file_path, data="", is_file=True) # Example adjustment
                translation_logger.info(f"Using updated file path: {updated_file_path}")
        else:
                file_path_to_process = file_path
                temp_file_obj = input_file # Use original object

        current_end_row = end_row
        if row_selection == "所有行":
            try:
                current_end_row = FileReaderFactory.count_rows(file_path_to_process)
                translation_logger.info(f"Row selection 'all rows' - counted {current_end_row} rows in {file_name}")
            except Exception as e:
                error_msg = f"Could not count rows for {file_name}: {e}. Skipping file or using default end_row."
                error_logger.error(error_msg)
                print(error_msg)
                continue # Or handle error differently

        try:
            translation_logger.info(f"Starting translation for {file_name} - rows {start_row} to {current_end_row}")
            process_time, output_file = translate_excel(temp_file_obj, start_row, current_end_row, start_column, target_column,
                                                        selected_model, selected_lora_model, selected_gpu, batch_size,
                                                        original_language, target_languages)
            if output_file and os.path.exists(output_file):
                    processed_file_path = os.path.join(processed_folder, os.path.basename(output_file))
                    shutil.move(output_file, processed_file_path)
                    processed_files.append(processed_file_path)
                    translation_logger.info(f"Successfully processed {file_name} -> {os.path.basename(processed_file_path)}")
            else:
                    error_msg = f"Translation failed or output file path invalid for {file_name}"
                    error_logger.error(error_msg)
                    print(error_msg)

        except Exception as e:
            error_msg = f"Error processing file {file_name}: {e}"
            error_logger.error(error_msg)
            error_logger.error(f"Traceback: {traceback.format_exc()}")
            print(error_msg)
            continue # Skip to next file on error
    zip_filename = os.path.join(folder_path, "processed_files.zip")
    if not processed_files:
            error_msg = "No files were processed successfully."
            error_logger.error(error_msg)
            print(error_msg)
            return "No files processed successfully.", None # Return None for zip file

    translation_logger.info(f"Creating zip file with {len(processed_files)} processed files")
    try:
        with zipfile.ZipFile(zip_filename, 'w') as zipf:
            for file in processed_files:
                if os.path.exists(file):
                        zipf.write(file, os.path.basename(file))
                        translation_logger.info(f"Added {os.path.basename(file)} to zip")
                        print(f"File {file} added to zip.")
                else:
                        warning_msg = f"Warning: Processed file {file} not found for zipping."
                        app_logger.warning(warning_msg)
                        print(warning_msg)
    except Exception as e:
        error_msg = f"Error creating zip file {zip_filename}: {e}"
        error_logger.error(error_msg)
        error_logger.error(f"Traceback: {traceback.format_exc()}")
        print(error_msg)
        return f"Error creating zip file: {e}", None

    end_time = time.time()
    total_time = int(end_time - start_time)
    translation_logger.info(f"Excel folder translation completed - Total time: {total_time}s, Files processed: {len(processed_files)}")
    print(f"Total process time: {total_time}s")
    print(f"Processed files added to zip: {processed_files}")
    return f"Total process time: {total_time}s. {len(processed_files)} file(s) processed.", zip_filename
    

@log_function_call('app')
def word_to_markdown(docx_path, output_dir="images"):
    """
    将指定的 .docx 文件转换为 Markdown 格式，并提取其中的图片。
    """
    app_logger.info(f"Converting Word document to Markdown: {os.path.basename(docx_path)}")
    app_logger.info(f"Output directory for images: {output_dir}")

    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
        app_logger.info(f"Created output directory: {output_dir}")
    def load_from_xml_v2(baseURI, rels_item_xml):
        srels = _SerializedRelationships()
        if rels_item_xml is not None:
            try:
                rels_elm = parse_xml(rels_item_xml)
                for rel_elm in rels_elm.Relationship_lst:
                    if hasattr(rel_elm, 'target_ref') and rel_elm.target_ref not in ('../NULL', 'NULL', None):
                        srels._srels.append(_SerializedRelationship(baseURI, rel_elm))
            except Exception as e:
                print(f"Error parsing relationships XML: {e}")
        return srels
    _SerializedRelationships.load_from_xml = load_from_xml_v2
    # --- Monkey Patching 结束 ---

    def iter_block_items(parent):
        if hasattr(parent, '_element') and hasattr(parent._element, 'body'):
            parent_elm = parent._element.body
            for child in parent_elm.iterchildren():
                if child.tag == qn('w:p'):
                    yield Paragraph(child, parent)
                elif child.tag == qn('w:tbl'):
                    yield Table(child, parent)

    try:
        doc = Document(docx_path)
    except Exception as e:
        print(f"Error opening document {docx_path}: {e}")
        return ""

    md_content = ""
    image_counter = 1
    image_paths_generated = []

    for block in iter_block_items(doc):
        if isinstance(block, Paragraph):
            para = block # para 是一个 Paragraph 对象
            if para._element.xpath('.//w:drawing'):
                for rId in para._element.xpath(".//a:blip/@r:embed"):
                    if rId and hasattr(doc.part, 'related_parts') and rId in doc.part.related_parts:
                        try:
                            image_part = doc.part.related_parts[rId]
                            image_bytes = image_part.blob
                            ext = os.path.splitext(image_part.partname)[-1].lower() or ".png"
                            if ext not in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff']:
                                ext = ".png"
                            
                            image_name = f"image_{image_counter}{ext}"
                            image_path = os.path.join(output_dir, image_name)
                            
                            with open(image_path, 'wb') as f:
                                f.write(image_bytes)
                            
                            md_image_path = image_name
                            md_content += f"![{image_name}]({md_image_path})\n\n"
                            image_paths_generated.append(image_path)
                            image_counter += 1
                        except Exception as e:
                            print(f"Error processing image resource {rId}: {e}")
            para_text = para.text.strip()
            if para_text:
                if para.style and para.style.name.startswith('Heading'):
                    try:
                        level = int(para.style.name.split()[-1])
                        md_content += f"{'#' * level} {para_text}\n\n"
                    except (ValueError, IndexError):
                        md_content += f"### {para_text}\n\n"
                else:
                    md_content += f"{para_text}\n\n"

        elif isinstance(block, Table):
            table = block
            md_content += "\n"
            rows = table.rows
            if len(rows) > 0:
                header_cells = rows[0].cells
                header = "| " + " | ".join((cell.text or "").strip().replace('\n', ' ') for cell in header_cells) + " |\n"
                md_content += header
                md_content += "| " + " | ".join(['---'] * len(header_cells)) + " |\n"
                for row in rows[1:]:
                    row_cells = row.cells
                    row_text = "| " + " | ".join((cell.text or "").strip().replace('\n', ' ') for cell in row_cells) + " |\n"
                    md_content += row_text
            md_content += "\n"

    app_logger.info(f"Word to Markdown conversion completed - Generated {len(image_paths_generated)} images")
    app_logger.info(f"Generated images: {image_paths_generated}")
    print(f"Generated images: {image_paths_generated}")
    return md_content
@log_function_call('app')
def markdown_to_word(md_content, word_path, image_base_dir="images"):
    app_logger.info(f"Converting Markdown to Word document: {os.path.basename(word_path)}")
    app_logger.info(f"Image base directory: {image_base_dir}")

    md_content = md_content.replace('<', '&lt;').replace('>', '&gt;')
    try:
        html = markdown.markdown(md_content, extensions=['tables', 'fenced_code'])
        app_logger.info("Successfully converted Markdown to HTML")
    except Exception as e:
        error_msg = f"Error converting Markdown to HTML: {e}"
        error_logger.error(error_msg)
        error_logger.error(f"Traceback: {traceback.format_exc()}")
        print(error_msg)
        return

    try:
        soup = BeautifulSoup(html, 'html.parser')
        app_logger.info("Successfully parsed HTML with BeautifulSoup")
    except Exception as e:
        error_msg = f"Error parsing generated HTML: {e}"
        error_logger.error(error_msg)
        error_logger.error(f"Traceback: {traceback.format_exc()}")
        print(error_msg)
        return

    doc = Document()

    for element in soup.children:
        if not hasattr(element, 'name'):
            text = str(element).strip()
            if text:
                doc.add_paragraph(text)
            continue

        if element.name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
            try:
                level = int(element.name[1])
                doc.add_heading(element.get_text(strip=True), level=level)
            except (ValueError, IndexError):
                doc.add_heading(element.get_text(strip=True), level=3)
        
        elif element.name == 'p':
            if element.find('img'):
                img_tag = element.find('img')
                img_src = img_tag.get('src')
                alt_text = img_tag.get('alt', '')
                img_path = os.path.join(image_base_dir, img_src)
                if os.path.exists(img_path):
                    try:
                        doc.add_picture(img_path, width=Inches(5.5))
                    except Exception as e:
                        print(f"Warning: Could not add picture {img_path}. Error: {e}")
                        doc.add_paragraph(f"[Image: {alt_text or img_src}]")
                else:
                    print(f"Warning: Could not find image file at a constructed path: {img_path} (source was: {img_src})")
                    doc.add_paragraph(f"[Image not found: {alt_text or img_src}]")
            else:
                text = element.get_text()
                if text.strip():
                    doc.add_paragraph(text)
        
        elif element.name == 'ul':
            for li in element.find_all('li', recursive=False):
                doc.add_paragraph(li.get_text(strip=True), style='List Bullet')
        elif element.name == 'ol':
            for li in element.find_all('li', recursive=False):
                doc.add_paragraph(li.get_text(strip=True), style='List Number')
        elif element.name == 'pre' and element.find('code'):
            code_text = element.find('code').get_text()
            p = doc.add_paragraph(style='BodyText')
            run = p.add_run(code_text)
        elif element.name == 'table':
            rows = element.find_all('tr')
            if not rows: continue

            num_cols = len(rows[0].find_all(['th', 'td']))
            if num_cols == 0: continue
            
            table = doc.add_table(rows=1, cols=num_cols)
            table.style = 'Table Grid'
            
            header_cells = rows[0].find_all(['th', 'td'])
            for i, cell in enumerate(header_cells):
                table.cell(0, i).text = cell.get_text(strip=True)
            
            for row_data in rows[1:]:
                row_cells_data = row_data.find_all('td')
                new_row = table.add_row().cells
                for i, cell in enumerate(row_cells_data):
                    if i < num_cols:
                        new_row[i].text = cell.get_text(strip=True)
    try:
        doc.save(word_path)
        app_logger.info(f"Successfully saved Word document: {os.path.basename(word_path)}")
    except Exception as e:
        error_msg = f"Error saving Word document to {word_path}: {e}"
        error_logger.error(error_msg)
        error_logger.error(f"Traceback: {traceback.format_exc()}")
        print(error_msg)

def update_lora_and_explanation(selected_model):
    """当模型改变时，只更新Lora模型列表和模型介绍。"""
    model_path = available_models.get(selected_model)
    lora_list = ['']
    model_explanation = "模型路径未找到或README.md缺失。"

    if model_path:
        readme_path = os.path.join(model_path, 'README.md')
        if os.path.isfile(readme_path):
            try:
                with open(readme_path, 'r', encoding='utf-8') as file:
                    model_explanation = file.read()
            except Exception as e:
                model_explanation = f"读取README.md时出错: {e}"
        try:
            lora_list.extend([f for f in os.listdir(model_path) if
                                os.path.isdir(os.path.join(model_path, f)) and not f.startswith('.') and not f.startswith('_')])
        except Exception as e:
            print(f"列出Lora模型时出错 {model_path}: {e}")
            
    return gr.Dropdown(choices=lora_list, value=''), gr.Textbox(value=model_explanation)

@log_function_call('translation')
def translate_excel_fixed_languages(input_file, selected_model = 'gpt-4.1-mini', max_workers=10) -> Tuple[str, Optional[str]]:
    """
    使用 OpenAI API 并行翻译 Excel 文件中的指定列。

    该函数严格遵循固定的列名（简体中文(源), English, 繁體中文）作为翻译参考，
    并将结果填充到 TARGET_COLUMNS 定义的各语言列中。

    Args:
        input_file: Gradio UI 上传的文件对象。
        max_workers (int): 并发处理的线程数。

    Returns:
        A tuple containing:
        - str: 处理过程和结果的状态信息。
        - Optional[str]: 处理成功后输出文件的路径，失败则为 None。
    """
    SIMPLE_COLUMN_NAME = "简体中文(源)"
    ENGLISH_COLUMN_NAME = "English"
    TRANS_COLUMN_NAME = "繁體中文"

    # 需要翻译的目标语言列名列表
    TARGET_COLUMNS = [
        "印尼语", "匈牙利语", "葡萄牙语", "泰语", "土耳其语", "越南语", "俄语",
        "阿拉伯语", "芬兰语", "丹麦语", "荷兰语", "波兰语", "法语", "德语",
        "日语", "挪威语", "希伯来语", "韩语", "西班牙语", "捷克语", "意大利语",
        "瑞典语", "希腊语", "马来语", "斯洛伐克语", "柬埔寨语", "罗马尼亚语",
        "克罗地亚语", "乌兹别克语", "缅甸语"
    ]

    if not input_file:
        error_msg = "错误：请先上传一个Excel文件。"
        error_logger.error(error_msg)
        return error_msg, None

    file_path = input_file.name
    translation_logger.info(f"Starting fixed languages translation - File: {os.path.basename(file_path)}")
    translation_logger.info(f"Model: {selected_model}, Max workers: {max_workers}")
    translation_logger.info(f"Target languages count: {len(TARGET_COLUMNS)}")

    client = OpenAI()

    start_time = time.time()
    status_messages = [f"▶ 开始处理文件: {os.path.basename(file_path)}"]

    try:
        df = pd.read_excel(file_path)
        status_messages.append(f"✔ 成功读取 Excel 文件，共 {len(df)} 行数据。")
        required_columns = [SIMPLE_COLUMN_NAME, ENGLISH_COLUMN_NAME, TRANS_COLUMN_NAME]
        if not all(col in df.columns for col in required_columns):
            missing_cols = [col for col in required_columns if col not in df.columns]
            return f"错误：输入文件缺少必需的列: {', '.join(missing_cols)}。请检查文件格式。", None
        for col in TARGET_COLUMNS:
            if col not in df.columns:
                df[col] = None
        all_target_languages = [TRANS_COLUMN_NAME] + TARGET_COLUMNS
        total_languages = len(all_target_languages)

        for lang_idx, target_lang_column in enumerate(all_target_languages):
            status_messages.append(f"\n--- ({lang_idx + 1}/{total_languages}) 正在处理: {target_lang_column} ---")
            print(f"\n--- Processing: {target_lang_column} ---")
            translation_cache = {}
            def generate_translation(index, row_data):
                """为单行数据生成翻译的核心函数"""
                if pd.notna(row_data.get(target_lang_column)):
                    return index, None, None # 返回 None 表示无需更新
                english_text = str(row_data[ENGLISH_COLUMN_NAME])
                if english_text in translation_cache:
                    return index, translation_cache[english_text], "cache"
                if target_lang_column == TRANS_COLUMN_NAME:
                    simplified_text = str(row_data[SIMPLE_COLUMN_NAME])
                    translated_text = zhconv.convert(simplified_text, 'zh-tw')
                    translation_cache[english_text] = translated_text
                    return index, translated_text, "zhconv"
                try:
                    completion = client.chat.completions.create(
                        model=selected_model,
                        messages=[
                            {"role": "user", "content": f"Translate the following sentence or word from English to {SIMPLE_COLUMN_NAME}: {english_text}, please directly translate it and do not output any extra content"},
                            {"role": "assistant", "content": str(row_data[SIMPLE_COLUMN_NAME])},
                            {"role": "user", "content": f"Translate the following sentence or word from English to {TRANS_COLUMN_NAME}: {english_text}, please directly translate it and do not output any extra content"},
                            {"role": "assistant", "content": zhconv.convert(str(row_data[SIMPLE_COLUMN_NAME]), 'zh-tw')},
                            {"role": "user", "content": f"Translate the following sentence or word from English to {target_lang_column}: {english_text}, please directly translate it and do not output any extra content"}
                        ],
                        temperature=0.0,
                        max_tokens=200
                    )
                    translated_text = completion.choices[0].message.content.strip()
                    translation_cache[english_text] = translated_text
                    return index, translated_text, "api"
                except Exception as api_error:
                    return index, f"API_ERROR: {api_error}", "error"
            tasks_to_process = [(index, row) for index, row in df.iterrows() if pd.isna(row.get(target_lang_column))]
            if not tasks_to_process:
                status_messages.append(f"✔ '{target_lang_column}' 列已全部翻译，跳过。")
                print(f"'{target_lang_column}' column is already fully translated. Skipping.")
                continue

            status_messages.append(f"找到 {len(tasks_to_process)} 个待翻译条目，开始处理...")
            
            with ThreadPoolExecutor(max_workers=max_workers) as executor:
                future_to_index = {executor.submit(generate_translation, index, row): index for index, row in tasks_to_process}
                for future in as_completed(future_to_index):
                    index, result, source = future.result()
                    if result is not None and source != "error":
                        df.at[index, target_lang_column] = result
                    elif source == "error":
                        print(f"Error processing row {index}: {result}")

        processed_dir = os.path.join(os.path.dirname(file_path), 'processed_openai')
        os.makedirs(processed_dir, exist_ok=True)
        base_name = os.path.basename(file_path)
        name, ext = os.path.splitext(base_name)
        timestamp = time.strftime("%Y%m%d_%H%M%S")
        output_file_path = os.path.join(processed_dir, f"{name}_translated_{timestamp}{ext}")
        
        df.to_excel(output_file_path, index=False)

    except Exception as e:
        error_message = f"处理Excel文件时发生严重错误: {e}"
        error_logger.error(f"Critical error in fixed languages translation: {e}")
        error_logger.error(f"Traceback: {traceback.format_exc()}")
        status_messages.append(f"\n❌ {error_message}")
        print(error_message)
        return "\n".join(status_messages), None

    end_time = time.time()
    total_time = int(end_time - start_time)
    translation_logger.info(f"Fixed languages translation completed successfully - Total time: {total_time}s")
    translation_logger.info(f"Output file: {os.path.basename(output_file_path)}")

    status_messages.append(f"\n🎉 所有翻译任务完成！总耗时: {total_time}秒。")
    status_messages.append(f"✔ 结果已保存至: {os.path.basename(output_file_path)}")

    return "\n".join(status_messages), output_file_path

def extract_complex_blocks(md_content: str):
    """
    使用占位符提取Markdown中的复杂块（图片、表格、代码块）。(此函数不变)
    """
    pattern = re.compile(
        r"(!\[.*?\]\(.*?\))|"
        r"((?:\|.*\|[\r\n]+)+(?:\|-+\|.*[\r\n]+)+(?:\|.*\|[\r\n]?)*)|"
        r"(```[\s\S]*?```)"
        , re.MULTILINE)
    blocks = {}
    def replacer(match):
        placeholder = f"__COMPLEX_BLOCK_{len(blocks)}__"
        blocks[placeholder] = match.group(0)
        return placeholder
    clean_md = pattern.sub(replacer, md_content)
    return clean_md, blocks

def restore_complex_blocks(translated_content: str, blocks: dict) -> str:
    """
    将占位符替换回其原始的复杂块内容。(此函数不变)
    """
    for placeholder, original_block in blocks.items():
        translated_content = translated_content.replace(placeholder, original_block)
    return translated_content


@log_function_call('translation')
def translate_markdown_folder(translating_files: list[NamedString],
                            selected_model: Optional[str], selected_lora_model: Optional[str],
                            selected_gpu: Optional[str], batch_size: int,
                            original_language: Optional[str], target_language: Optional[str]):
    start_time = time.time()

    if not translating_files:
        error_msg = "No files uploaded for markdown translation"
        error_logger.error(error_msg)
        return "No files uploaded", None

    folder_path = os.path.dirname(translating_files[0].name)
    translation_logger.info(f"Starting markdown folder translation - {len(translating_files)} files")
    translation_logger.info(f"Folder path: {folder_path}")
    translation_logger.info(f"Model: {selected_model}, GPU: {selected_gpu}, Batch size: {batch_size}")
    translation_logger.info(f"Languages - From: {original_language}, To: {target_language}")

    processed_files = []
    temp_image_dir = os.path.join(folder_path, "temp_images_from_docx")

    processed_folder = os.path.join(folder_path, 'processed')
    os.makedirs(processed_folder, exist_ok=True)
    translation_logger.info(f"Created processed folder: {processed_folder}")

    if os.path.exists(temp_image_dir):
        shutil.rmtree(temp_image_dir)
    os.makedirs(temp_image_dir, exist_ok=True)
    translation_logger.info(f"Created temp image directory: {temp_image_dir}")

    for input_file in translating_files:
        file_path = input_file.name
        file_name, file_ext = os.path.splitext(os.path.basename(file_path))
        output_file_path = None

        translation_logger.info(f"Processing file: {os.path.basename(file_path)} (type: {file_ext})")

        try:
            if file_ext.lower() == '.pptx':
                translation_logger.info(f"Processing PowerPoint file: {file_name}")
                def extract_text_from_shape(shape, run_list, text_list):
                    """递归提取所有文本，包括文本框、表格和嵌套形状"""
                    if hasattr(shape, "text_frame") and shape.text_frame is not None:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run_list.append(run)
                                text_list.append(run.text)
                    elif getattr(shape, "has_table", False):
                        table = shape.table
                        for row in table.rows:
                            for cell in row.cells:
                                if cell.text_frame is not None:
                                    for paragraph in cell.text_frame.paragraphs:
                                        for run in paragraph.runs:
                                            run_list.append(run)
                                            text_list.append(run.text)
                    elif hasattr(shape, "shapes"):
                        for sub_shape in shape.shapes:
                            extract_text_from_shape(sub_shape, run_list, text_list)

                prs = Presentation(file_path)
                run_list = []
                text_list = []

                for slide in prs.slides:
                    for shape in slide.shapes:
                        extract_text_from_shape(shape, run_list, text_list)  # 确保提取所有文本
                translated_segments = translate(text_list, selected_model, selected_lora_model, selected_gpu,
                                                batch_size, original_language, target_language)
                for run, translated in zip(run_list, translated_segments):
                    run.text = " " + translated[0]["generated_translation"]
                output_file_path = os.path.join(processed_folder, os.path.basename(file_name + '.pptx'))
                prs.save(output_file_path)
                processed_files.append(output_file_path)

            elif file_ext.lower() in ['.xlsx', '.xls']:
                translation_logger.info(f"Processing Excel file: {file_name}")
                print(f"Processing Excel file: {file_path}")

                workbook = openpyxl.load_workbook(file_path)
                texts_to_translate = []
                cell_locations = []
                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    for row_idx, row in enumerate(sheet.iter_rows()):
                        for col_idx, cell in enumerate(row):
                            if cell.value:  # 只要单元格不为空
                                texts_to_translate.append(str(cell.value))
                                cell_locations.append((sheet_name, cell.row, cell.column))

                translation_logger.info(f"Extracted {len(texts_to_translate)} cells for translation from Excel file")

                if texts_to_translate:
                    target_lang_list = [target_language] if isinstance(target_language, str) else target_language
                    translated_results = translate(texts_to_translate, selected_model, selected_lora_model, selected_gpu,
                                                batch_size, original_language, target_lang_list)
                    translation_logger.info(f"Completed translation for Excel file - {len(translated_results)} results")

                    for i, location in enumerate(cell_locations):
                        translated_text = translated_results[i][0]['generated_translation']
                        sheet_name, row, col = location
                        workbook[sheet_name].cell(row=row, column=col, value=translated_text)

                output_filename_base = os.path.basename(file_name + '_translated')
                output_file_path = os.path.join(processed_folder, output_filename_base + file_ext)
                workbook.save(output_file_path)
                translation_logger.info(f"Saved translated Excel file: {output_filename_base + file_ext}")

            elif file_ext.lower() in ['.docx', '.md']:
                translation_logger.info(f"Processing document file: {file_name} (type: {file_ext})")
                md_content = ""
                file_is_word = False
                if file_ext.lower() == '.docx':
                    file_is_word = True
                    translation_logger.info(f"Converting Word document to markdown: {file_name}")
                    md_content = word_to_markdown(file_path, output_dir=temp_image_dir)
                elif file_ext.lower() == '.md':
                    translation_logger.info(f"Reading markdown file: {file_name}")
                    with open(file_path, 'r', encoding='utf-8') as f:
                        md_content = f.read()

                clean_md, protected_blocks = extract_complex_blocks(md_content)
                text_to_translate = [p for p in clean_md.split('\n\n')]
                translated_content = clean_md

                translation_logger.info(f"Extracted {len(text_to_translate)} paragraphs for translation")
                translation_logger.info(f"Protected {len(protected_blocks)} complex blocks from translation")

                if text_to_translate:
                    target_lang_list = [target_language] if isinstance(target_language, str) else target_language
                    translated_results = translate(text_to_translate, selected_model, selected_lora_model, selected_gpu,
                                                batch_size, original_language, target_lang_list)
                    translation_logger.info(f"Completed translation for document - {len(translated_results)} results")

                    translation_map = {
                        original: result_list[0]['generated_translation']
                        for original, result_list in zip(text_to_translate, translated_results)
                    }

                    temp_translated_content = []
                    for para in clean_md.split('\n\n'):
                        temp_translated_content.append(translation_map.get(para, para))
                    translated_content = '\n\n'.join(temp_translated_content)

                final_md_content = restore_complex_blocks(translated_content, protected_blocks)

                output_filename_base = os.path.basename(file_name + '_translated')
                if file_is_word:
                    output_file_path = os.path.join(processed_folder, output_filename_base + '.docx')
                    translation_logger.info(f"Converting translated content back to Word document")
                    markdown_to_word(final_md_content, output_file_path, image_base_dir=temp_image_dir)
                else:
                    output_file_path = os.path.join(processed_folder, output_filename_base + '.md')
                    translation_logger.info(f"Saving translated markdown file")
                    with open(output_file_path, 'w', encoding='utf-8') as f:
                        f.write(final_md_content)
            else:
                warning_msg = f"Skipping unsupported file type: {file_path}"
                app_logger.warning(warning_msg)
                print(warning_msg)
                continue

            if output_file_path and os.path.exists(output_file_path):
                processed_files.append(output_file_path)
                translation_logger.info(f"Successfully processed file: {os.path.basename(output_file_path)}")

        except Exception as e:
            error_msg = f"CRITICAL ERROR processing file {file_path}: {e}"
            error_logger.error(error_msg)
            error_logger.error(f"Traceback: {traceback.format_exc()}")
            print(error_msg)
            import traceback
            traceback.print_exc()
            continue
    if not processed_files:
        error_msg = "No files processed successfully."
        error_logger.error(error_msg)
        if os.path.exists(temp_image_dir):
            shutil.rmtree(temp_image_dir)
        return error_msg, None

    translation_logger.info(f"Processing completed - {len(processed_files)} files processed successfully")

    output_path = None
    if len(processed_files) == 1:
        output_path = processed_files[0]
        translation_logger.info(f"Single file output: {os.path.basename(output_path)}")
    else:
        zip_filename = os.path.join(folder_path, "processed_files.zip")
        translation_logger.info(f"Creating zip file with {len(processed_files)} files")
        try:
            with zipfile.ZipFile(zip_filename, 'w') as zipf:
                for file in processed_files:
                    if os.path.exists(file):
                        zipf.write(file, os.path.basename(file))
                        translation_logger.info(f"Added {os.path.basename(file)} to zip")
            output_path = zip_filename
            translation_logger.info(f"Zip file created successfully: {os.path.basename(zip_filename)}")
        except Exception as e:
            error_msg = f"Error creating zip file: {e}"
            error_logger.error(error_msg)
            error_logger.error(f"Traceback: {traceback.format_exc()}")
            if os.path.exists(temp_image_dir):
                shutil.rmtree(temp_image_dir)
            return error_msg, None

    if os.path.exists(temp_image_dir):
        shutil.rmtree(temp_image_dir)
        translation_logger.info("Cleaned up temporary image directory")

    end_time = time.time()
    duration = int(end_time - start_time)
    translation_logger.info(f"Markdown folder translation completed - Total time: {duration}s, Files processed: {len(processed_files)}")
    return f"Total process time: {duration}s. {len(processed_files)} file(s) processed.", output_path
@log_function_call('app')
def glossary_check(input_folder, start_row, end_row, original_column, reference_column, translated_column,
                    row_selection, remark_column) -> tuple[str, Optional[str]]: # Return tuple[status, filepath]
    def contains_special_string(sentence):
        patterns = {
            "Content within <% ... %> should not be translated": r"<%.*?%>",  # Match <% ... %>
            "Special symbol %s should be contained": r"%s",  # Match %s
            "Special symbols {0}, {1}, {2}, etc., should not be translated": r"{\d+}",  # Match {0}, {1}, {2}, etc.
            "Special symbol %d should be contained": r"%d",  # Match %d
            "String {counts} should not be translated": r"{counts}",  # Match {counts}
            "String (value) should not be translated": r"\(value\)",  # Match (value)
            "String (Value) should not be translated": r"\(Value\)",  # Match (Value)
            "String (text) should not be translated": r"\(text\)",  # Match (text)
            "String (Text) should not be translated": r"\(Text\)",  # Match (Text)
            "String (message) should not be translated": r"\(message\)",  # Match (message)
            "String (Message) should not be translated": r"\(Message\)",  # Match (Message)
            "String (group) should not be translated": r"\(group\)",  # Match (group)
            "String (Group) should not be translated": r"\(Group\)",  # Match (Group)
            "Content within &{...}& should not be translated": r"&{.*?}&",  # Match &{...}&
            "String {} should be contained": r"{}",  # Match {}
            "Content within #...# should not be translated": r"#.*?#",  # Match #...#
            "Content within {{...}} should not be translated": r"{{.*?}}",  # Match {{...}}
            "Consecutive uppercase letters (AR, AP, SKU) should be contained": r"\b[A-Z]{2,}\b", # Use \b for word boundaries
            "CamelCase words (e.g., ServiceCode, LocStudio) should be contained": r"\b(?:[A-Z][a-z]+){2,}\b", # Use \b
            "Full links http:// should be contained": r"http://",  # Match full links containing "http://"
            "Full links https:// should be contained": r"https://",  # Match full links containing "https://"
            "Full file paths E:\\, D:\\, C:\\ should be contained": r"\b[CDE]:\\", # Use \b
            "Formula-like strings such as datediff(.*?,.*?,.*?) should not be translated": r"datediff\(.*?,.*?,.*?\)",
            "Strings like @BusinessFunction. ... @ should not be translated": r"@业务函数\..*?@",
            "CamelCase words starting with a lowercase letter (e.g., serviceCode, locStudio) should not be translated": r"\b[a-z]+[A-Z][a-zA-Z]*\b", # Use \b
            "String ${label} should not be translated": r"\$\{label\}",
            "String [${enum}] should not be translated": r"\[\$\{enum\}\]",
            "String ${max} should not be translated": r"\$\{max\}",
            "String ${min} should not be translated": r"\$\{min\}",
            "String ${len} should not be translated": r"\$\{len\}",
            "String ${pattern} should not be translated": r"\$\{pattern\}",
            "String [{{fievent}}] should not be translated": r"\[\{\{fievent\}\}\]",
            "String [{{accBook}}] should not be translated": r"\[\{\{accBook\}\}\]",
        }

        reasons = []
        matched_strings = []
        sentence_str = str(sentence) if sentence is not None else ""

        for reason, pattern in patterns.items():
            try:
                matches = re.findall(pattern, sentence_str)
                if matches:
                    if reason not in reasons:
                            reasons.append(reason)
                    for match in matches:
                            if match not in matched_strings:
                                matched_strings.append(match)
            except Exception as e:
                print(f"Regex error for pattern '{pattern}' on sentence '{sentence_str[:50]}...': {e}")


        return {
            "contains_special_string": bool(reasons),
            "reason": reasons,
            "matched_strings": matched_strings
        }

    app_logger.info(f"Starting glossary check - {len(input_folder) if input_folder else 0} files")
    app_logger.info(f"Parameters - Rows: {start_row}-{end_row}, Row selection: {row_selection}")
    app_logger.info(f"Columns - Original: {original_column}, Reference: {reference_column}, Translated: {translated_column}, Remark: {remark_column}")

    result = []
    excel_writer = ExcelFileWriter()
    processed_files = []
    output_zip_path = None # Initialize zip path

    if not input_folder:
            error_msg = "Error: No folder/files provided."
            error_logger.error(error_msg)
            return error_msg, None

    try:
            folder_path = os.path.dirname(input_folder[0].name) # Get dir from first file
            processed_folder = os.path.join(folder_path, 'processed_glossary_check')
            os.makedirs(processed_folder, exist_ok=True)
            app_logger.info(f"Created processed folder: {processed_folder}")
    except Exception as e:
            error_msg = f"Error creating processed folder: {e}"
            error_logger.error(error_msg)
            error_logger.error(f"Traceback: {traceback.format_exc()}")
            return error_msg, None


    for input_file in input_folder:
        file_path = input_file.name
        file_name_base = os.path.basename(file_path)
        file_name, file_ext = os.path.splitext(file_name_base)

        app_logger.info(f"Processing file for glossary check: {file_name_base}")

        if file_ext.lower() == '.xlsx':
            current_end_row = end_row
            try:
                if row_selection == "所有行":
                    current_end_row = FileReaderFactory.count_rows(file_path)
                    app_logger.info(f"Row selection 'all rows' - counted {current_end_row} rows in {file_name_base}")

                reader, fp = FileReaderFactory.create_reader(file_path) # fp might be None or file pointer
                original_inputs = reader.extract_text(file_path, original_column, start_row, current_end_row)
                reference_inputs = reader.extract_text(file_path, reference_column, start_row, current_end_row)
                translated_inputs = reader.extract_text(file_path, translated_column, start_row, current_end_row)

                app_logger.info(f"Extracted data from {file_name_base} - Original: {len(original_inputs)}, Reference: {len(reference_inputs)}, Translated: {len(translated_inputs)}")

                if fp: # Close file pointer if factory returned one
                        fp.close()

            except Exception as e:
                error_msg = f"Error reading {file_name_base}: {e}"
                error_logger.error(error_msg)
                error_logger.error(f"Traceback: {traceback.format_exc()}")
                result.append(error_msg)
                continue # Skip to next file

            result.append(f"Checking {file_name_base}:")
            app_logger.info(f"Starting glossary check for {file_name_base}")
            outputs = [] # Remarks to write back to Excel
            max_len = max(len(original_inputs), len(reference_inputs), len(translated_inputs))
            min_len = min(len(original_inputs), len(reference_inputs), len(translated_inputs))
            if max_len != min_len:
                    result.append(f"\tWarning: Column lengths differ ({len(original_inputs)}, {len(reference_inputs)}, {len(translated_inputs)}). Processing up to shortest length: {min_len}")

            for index in range(min_len):
                original_input = original_inputs[index]
                reference_input = reference_inputs[index]
                translated_input = translated_inputs[index]
                remark = ""
                original_str = str(original_input) if original_input is not None else ""
                reference_str = str(reference_input) if reference_input is not None else ""
                translated_str = str(translated_input) if translated_input is not None else ""
                special_check_result = contains_special_string(original_str)

                if special_check_result["contains_special_string"]:
                    missed_matches_info = [] # Store tuples of (missed_string, reason)
                    found_in_translation = True # Assume found initially

                    for matched_string in special_check_result["matched_strings"]:
                        if matched_string not in translated_str:
                            if matched_string not in reference_str:
                                pass
                            else:
                                found_in_translation = False
                                associated_reasons = [r for r, p in contains_special_string(original_str)["reason"].items() if re.search(p, matched_string)]
                                reason_text = associated_reasons[0] if associated_reasons else "Unknown reason" # Get first reason
                                missed_matches_info.append((matched_string, reason_text))


                    if not found_in_translation:
                        missed_items_str = ', '.join([f"'{info[0]}'" for info in missed_matches_info])
                        reasons_str = ', '.join(set([info[1] for info in missed_matches_info])) # Unique reasons
                        remark += f"MISSED: {missed_items_str}; REASON: {reasons_str}"
                        result.append(
                            f"\tROW: {start_row + index}, MISSED: {missed_items_str}, REASON: {reasons_str}")
                    else:
                            remark += "SPECIAL_STRINGS_OK"
                try:
                    if 'tokenizer' in globals() and tokenizer:
                            token_count = len(tokenizer.encode(original_str))
                            if token_count > 40: # Your threshold
                                length_remark = "EXCEEDS_40_TOKENS"
                                if remark: # Append to existing remark
                                    remark += f"; {length_remark}"
                                else:
                                    remark = length_remark
                                result.append(f"\tROW: {start_row + index}, {length_remark} ({token_count} tokens)")
                    else:
                        print("Warning: Tokenizer not loaded, skipping length check.")
                except Exception as e:
                    print(f"Error during tokenization for row {start_row + index}: {e}")
                outputs.append(remark if remark else "") # Append remark or empty string
            if len(outputs) != min_len:
                print(f"Warning: Length mismatch in remarks generation for {file_name_base}. Expected {min_len}, got {len(outputs)}.")
                outputs = outputs[:min_len] + [""] * (min_len - len(outputs))
            try:
                output_file = excel_writer.write_list(file_path, outputs, remark_column, start_row, current_end_row)
                processed_file_path = os.path.join(processed_folder, os.path.basename(output_file))
                shutil.move(output_file, processed_file_path)
                processed_files.append(processed_file_path)
                result.append(f"{file_name_base} check completed. Results saved to processed folder.")

            except Exception as e:
                result.append(f"Error writing remarks or moving file for {file_name_base}: {e}")
                if 'output_file' in locals() and os.path.exists(output_file):
                    try: os.remove(output_file)
                    except OSError: pass


        else:
            warning_msg = f"Skipping non-Excel file: {file_name_base}"
            app_logger.warning(warning_msg)
            result.append(warning_msg)

    if processed_files:
        zip_filename_base = "glossary_check_results.zip"
        output_zip_path = os.path.join(folder_path, zip_filename_base) # Save zip in original upload dir
        app_logger.info(f"Creating zip file with {len(processed_files)} processed files")
        try:
            with zipfile.ZipFile(output_zip_path, 'w') as zipf:
                for file in processed_files:
                    if os.path.exists(file):
                        zipf.write(file, os.path.basename(file))
                        app_logger.info(f"Added {os.path.basename(file)} to zip")
                        print(f"Adding {os.path.basename(file)} to zip.")
                    else:
                        warning_msg = f"Warning: File {file} not found for zipping."
                        app_logger.warning(warning_msg)
                        print(warning_msg)
            result.append(f"Processed files zipped to {zip_filename_base}")
            app_logger.info(f"Glossary check completed successfully - Zip file created: {zip_filename_base}")
        except Exception as e:
            error_msg = f"Error creating zip file: {e}"
            error_logger.error(error_msg)
            error_logger.error(f"Traceback: {traceback.format_exc()}")
            result.append(error_msg)
            output_zip_path = None # Indicate zip creation failed
    else:
        error_msg = "No files were processed successfully."
        error_logger.error(error_msg)
        result.append(error_msg)

    app_logger.info("Glossary check process completed")
    return "\n".join(result), output_zip_path # Return status string and path to zip (or None)


def webui():

    def update_choices(selected_model):
        app_logger.info(f"User selected model: {selected_model}")

        model_path = available_models.get(selected_model) # 使用 .get() 更安全
        original_language_choices = []
        target_language_choices = []
        lora_list = ['']
        model_explanation = "Model path not found or README.md missing."

        if model_path:
            app_logger.info(f"Loading model configuration from: {model_path}")
            support_language_path = os.path.join(model_path, 'support_language.json')
            readme_path = os.path.join(model_path, 'README.md')

            if os.path.isfile(readme_path):
                try:
                    with open(readme_path, 'r', encoding='utf-8') as file:
                        model_explanation = file.read()
                except Exception as e:
                    print(f"Error reading README.md: {e}")
                    model_explanation = f"Error reading README.md: {e}"

            try:
                with open(support_language_path, 'r') as file:
                    support_languages = json.load(file)
                    original_language_choices = support_languages.get("original_language", [])
                    target_language_choices = support_languages.get("target_language", [])
            except Exception as e:
                print(f"Error reading support_language.json: {e}")
                # Keep choices empty if file read fails

            try:
                lora_list = [''] + [f for f in os.listdir(model_path) if
                                    os.path.isdir(os.path.join(model_path, f)) and not f.startswith('.') and not f.startswith(
                                        '_')]
            except Exception as e:
                print(f"Error listing lora models in {model_path}: {e}")
                lora_list = [''] # Reset to default if error

        # --- 重要: 返回更新后的组件，同时确保默认值在 choices 中 ---
        # 如果默认值不在动态加载的choices里，Gradio 可能无法正确显示默认值
        # 这里假设 'Chinese' 和 'English' 通常会存在于加载的语言列表中
        # 如果模型不支持这些语言，默认值可能不会被选中
        return (gr.Dropdown(choices=original_language_choices, value=default_original_language if default_original_language in original_language_choices else None),
                gr.Dropdown(choices=target_language_choices, value=default_target_language_multi if default_target_language_single in target_language_choices else None, multiselect=True), # 处理多选的情况
                gr.Dropdown(choices=lora_list), # Lora 模型通常没有默认值，除了空字符串
                model_explanation)

    
    # --- Gradio UI Layout ---
    with gr.Blocks(title="yonyou translator") as interface:
        # Define initial choices based on the default model, if possible
        initial_original_choices = []
        initial_target_choices = []
        initial_lora_choices = ['']
        initial_explanation = "Select a model to see details."
        if default_model_name in available_models:
            try:
                 _model_path = available_models[default_model_name]
                 _sl_path = os.path.join(_model_path, 'support_language.json')
                 _readme_path = os.path.join(_model_path, 'README.md')

                 if os.path.exists(_readme_path):
                      with open(_readme_path, 'r', encoding='utf-8') as f: initial_explanation = f.read()
                 if os.path.exists(_sl_path):
                      with open(_sl_path, 'r') as f:
                           _langs = json.load(f)
                           initial_original_choices = _langs.get("original_language", [])
                           initial_target_choices = _langs.get("target_language", [])
                 try:
                    initial_lora_choices = [''] + [f for f in os.listdir(_model_path) if
                                     os.path.isdir(os.path.join(_model_path, f)) and not f.startswith('.') and not f.startswith('_')]
                 except Exception: pass # Ignore lora errors on initial load

            except Exception as e:
                 print(f"Error pre-loading defaults for {default_model_name}: {e}")
        with gr.Tabs():
            with gr.TabItem("Excel Translator"):
                with gr.Row():
                    with gr.Column():
                        input_file = gr.File()
                        with gr.Row():
                            start_row = gr.Number(value=2, label="起始行")
                            end_row = gr.Number(value=100001, label="终止行")
                            target_column = gr.Textbox(value="G", label="目标列")
                            start_column = gr.Textbox(value="H", label="结果写入列")
                        with gr.Row():
                            selected_model_excel = gr.Dropdown(choices=list(available_models.keys()), label="选择基模型", value=default_model_name)
                            selected_lora_model_excel = gr.Dropdown(choices=initial_lora_choices, label="选择Lora模型", value='')
                            selected_gpu_excel = gr.Dropdown(choices=available_gpus, label="选择GPU", value=available_gpus[0] if available_gpus else None)
                            batch_size_excel = gr.Number(value=10, label="批处理大小")
                        with gr.Row():
                            original_language_excel = gr.Dropdown(choices=initial_original_choices, label="原始语言", value=default_original_language)
                            target_languages_excel = gr.Dropdown(choices=initial_target_choices, label="目标语言", multiselect=True, value=default_target_language_multi)
                        translate_button_excel = gr.Button("Translate")
                    with gr.Column():
                        model_explanation_textbox_excel = gr.Textbox(label="模型介绍", lines=5, value=initial_explanation)
                        output_text_excel = gr.Textbox(label="输出文本")
                        output_file_excel = gr.File(label="翻译文件下载")

                # Link change event
                selected_model_excel.change(update_choices,
                                            inputs=[selected_model_excel],
                                            outputs=[original_language_excel, target_languages_excel, selected_lora_model_excel, model_explanation_textbox_excel])
                # Link click event
                translate_button_excel.click(translate_excel,
                                             inputs=[input_file, start_row, end_row, start_column, target_column,
                                                     selected_model_excel, selected_lora_model_excel, selected_gpu_excel, batch_size_excel,
                                                     original_language_excel, target_languages_excel],
                                             outputs=[output_text_excel, output_file_excel])

            with gr.TabItem("Text Translator"):
                 with gr.Row():
                     with gr.Column():
                         input_text_text = gr.Textbox(label="输入文本", lines=3)
                         with gr.Row():
                              selected_model_text = gr.Dropdown(choices=list(available_models.keys()), label="选择基模型", value=default_model_name)
                              selected_lora_model_text = gr.Dropdown(choices=initial_lora_choices, label="选择Lora模型", value='')
                              selected_gpu_text = gr.Dropdown(choices=available_gpus, label="选择GPU", value=available_gpus[0] if available_gpus else None)
                              batch_size_text = gr.Number(value=1, label="批处理大小", visible=False) # Usually 1 for single text
                         with gr.Row():
                              original_language_text = gr.Dropdown(choices=initial_original_choices, label="原始语言", value=default_original_language)
                              target_languages_text = gr.Dropdown(choices=initial_target_choices, label="目标语言", multiselect=True, value=default_target_language_multi)
                         translate_button_text = gr.Button("Translate")
                     with gr.Column():
                          model_explanation_textbox_text = gr.Textbox(label="模型介绍", lines=5, value=initial_explanation)
                          output_text_text = gr.Textbox(label="输出文本", lines=5)

                 selected_model_text.change(update_choices,
                                           inputs=[selected_model_text],
                                           outputs=[original_language_text, target_languages_text, selected_lora_model_text, model_explanation_textbox_text])
                 translate_button_text.click(translate,
                                            inputs=[input_text_text, selected_model_text, selected_lora_model_text, selected_gpu_text,
                                                    batch_size_text, original_language_text, target_languages_text],
                                            outputs=output_text_text) # Output only to text box

            with gr.TabItem("Folder Excel Translator"):
                with gr.Row():
                    with gr.Column():
                        input_folder_fexcel = gr.File(file_count="directory", label="选择Excel文件所在文件夹")
                        with gr.Row():
                            start_row_fexcel = gr.Number(value=yaml_data["excel_config"]["default_start_row"], label="起始行")
                        with gr.Row():
                            row_selection_fexcel = gr.Radio(choices=["特定行", "所有行"], label="行选择", value="特定行")
                            end_row_fexcel = gr.Number(value=yaml_data["excel_config"]["default_end_row"], label="终止行", visible=True)
                        row_selection_fexcel.change(update_row_selection, inputs=row_selection_fexcel, outputs=end_row_fexcel)
                        with gr.Row():
                            target_column_fexcel = gr.Textbox(value=yaml_data["excel_config"]["default_target_column"], label="目标列")
                            start_column_fexcel = gr.Textbox(value=yaml_data["excel_config"]["default_start_column"], label="结果写入列")
                        with gr.Row():
                            selected_model_fexcel = gr.Dropdown(choices=list(available_models.keys()), label="选择基模型", value=default_model_name)
                            selected_lora_model_fexcel = gr.Dropdown(choices=initial_lora_choices, label="选择Lora模型", value='')
                            selected_gpu_fexcel = gr.Dropdown(choices=available_gpus, label="选择GPU", value=available_gpus[0] if available_gpus else None)
                            batch_size_fexcel = gr.Number(value=yaml_data.get("excel_config", {}).get("default_batch_size", 10), label="批处理大小", visible=True) # Use default from yaml or 10
                        with gr.Row():
                            original_language_fexcel = gr.Dropdown(choices=initial_original_choices, label="原始语言", value=default_original_language)
                            target_languages_fexcel = gr.Dropdown(choices=initial_target_choices, label="目标语言", multiselect=True, value=default_target_language_multi)
                        translate_button_fexcel = gr.Button("Translate Folder")
                    with gr.Column():
                        model_explanation_textbox_fexcel = gr.Textbox(label="模型介绍", lines=5, value=initial_explanation)
                        output_text_fexcel = gr.Textbox(label="处理状态", lines=5)
                        output_folder_fexcel = gr.File(label="下载处理后的Zip文件")

                selected_model_fexcel.change(update_choices,
                                            inputs=[selected_model_fexcel],
                                            outputs=[original_language_fexcel, target_languages_fexcel, selected_lora_model_fexcel, model_explanation_textbox_fexcel])
                translate_button_fexcel.click(translate_excel_folder,
                                            inputs=[input_folder_fexcel, start_row_fexcel, end_row_fexcel, start_column_fexcel, target_column_fexcel,
                                                    selected_model_fexcel, selected_lora_model_fexcel, selected_gpu_fexcel, batch_size_fexcel,
                                                    original_language_fexcel, target_languages_fexcel, row_selection_fexcel],
                                            outputs=[output_text_fexcel, output_folder_fexcel])


            with gr.TabItem("Folder Markdown & Docx Translator"):
                with gr.Row():
                    with gr.Column():
                         # Allow uploading directory or individual files
                        input_folder_mdoc = gr.File(file_count="multiple", file_types=['.md', '.docx', '.pptx', '.xlsx', '.xls'], label="选择Markdown, Docx, PPTX文件或文件夹") # Allow multiple file types
                        with gr.Row():
                            selected_model_mdoc = gr.Dropdown(choices=list(available_models.keys()), label="选择基模型", value=default_model_name)
                            selected_lora_model_mdoc = gr.Dropdown(choices=initial_lora_choices, label="选择Lora模型", value='')
                            selected_gpu_mdoc = gr.Dropdown(choices=available_gpus, label="选择GPU", value=available_gpus[0] if available_gpus else None)
                            batch_size_mdoc = gr.Number(value=5, label="批处理大小", visible=True) # Adjust default batch size as needed
                        with gr.Row():
                            original_language_mdoc = gr.Dropdown(choices=initial_original_choices, label="原始语言", value=default_original_language)
                            # Changed target_language to single select for this tab as per original code structure
                            target_language_mdoc = gr.Dropdown(choices=initial_target_choices, label="目标语言", value=default_target_language_single)
                        translate_button_mdoc = gr.Button("Translate Folder/Files")
                    with gr.Column():
                        model_explanation_textbox_mdoc = gr.Textbox(label="模型介绍", lines=5, value=initial_explanation)
                        output_text_mdoc = gr.Textbox(label="处理状态", lines=5)
                        output_folder_mdoc = gr.File(label="下载处理后的Zip文件")

                selected_model_mdoc.change(update_choices,
                                            inputs=[selected_model_mdoc],
                                            outputs=[original_language_mdoc, target_language_mdoc, selected_lora_model_mdoc, model_explanation_textbox_mdoc]) # Map 2nd output to single dropdown
                translate_button_mdoc.click(translate_markdown_folder,
                                          inputs=[input_folder_mdoc, selected_model_mdoc, selected_lora_model_mdoc, selected_gpu_mdoc,
                                                  batch_size_mdoc, original_language_mdoc, target_language_mdoc], # Pass single target lang
                                          outputs=[output_text_mdoc, output_folder_mdoc])

            with gr.TabItem("33语翻译"):
                with gr.Row():
                    with gr.Column(scale=2):
                        input_file_fixed = gr.File(label="上传待翻译的Excel文件")
                        translate_button_fixed = gr.Button("开始翻译", variant="primary")
                    
                    with gr.Column(scale=1):
                        model_explanation_textbox_fixed = gr.Textbox(label="模型介绍", lines=10, value=initial_explanation, interactive=False)
                        output_text_fixed = gr.Textbox(label="处理状态与日志", lines=10, interactive=False)
                        output_file_fixed = gr.File(label="下载翻译后的文件")
                
                translate_button_fixed.click(
                    translate_excel_fixed_languages,
                    inputs=[
                        input_file_fixed
                    ],
                    outputs=[output_text_fixed, output_file_fixed]
                )

            with gr.TabItem("术语表校验"):
                 with gr.Row():
                     with gr.Column():
                         input_folder_gloss = gr.File(file_count="directory", label="选择包含待校验Excel文件的文件夹")
                         with gr.Row():
                             row_selection_gloss = gr.Radio(choices=["特定行", "所有行"], label="行选择", value="特定行")
                             start_row_gloss = gr.Number(value=yaml_data["excel_config"]["default_start_row"], label="起始行")
                             end_row_gloss = gr.Number(value=yaml_data["excel_config"]["default_end_row"], label="终止行", visible=True)
                         row_selection_gloss.change(update_row_selection, inputs=row_selection_gloss, outputs=end_row_gloss)
                         with gr.Row():
                              # Use defaults from YAML or hardcode if not present
                              default_orig_col = yaml_data.get("glossary_config", {}).get("original_column", "J")
                              default_ref_col = yaml_data.get("glossary_config", {}).get("reference_column", "G")
                              default_trans_col = yaml_data.get("glossary_config", {}).get("translated_column", "H")
                              default_remark_col = yaml_data.get("glossary_config", {}).get("remark_column", "I")
                              original_column_gloss = gr.Textbox(default_orig_col, label="原文列")
                              reference_column_gloss = gr.Textbox(default_ref_col, label="参考列")
                              translated_column_gloss = gr.Textbox(default_trans_col, label="已翻译列")
                              remark_column_gloss = gr.Textbox(default_remark_col, label="备注写入列")
                         glossary_check_button = gr.Button("开始检测")

                     with gr.Column():
                         output_text_gloss = gr.Textbox(label="检测结果摘要", lines=20, show_copy_button=True)
                         output_folder_gloss = gr.File(label="下载标注后的Zip文件")

                 glossary_check_button.click(glossary_check,
                                            inputs=[input_folder_gloss, start_row_gloss, end_row_gloss, original_column_gloss, reference_column_gloss,
                                                    translated_column_gloss, row_selection_gloss, remark_column_gloss],
                                            outputs=[output_text_gloss, output_folder_gloss])

    return interface


main_ui = webui()

if __name__ == "__main__":
    app_logger.info("Launching Gradio web interface")
    app_logger.info("Server configuration - Port: 8082, Host: 0.0.0.0, Share: False")

    try:
        # Consider adding server_name="0.0.0.0" to allow access from other devices on the network
        main_ui.launch(share=False, server_port=8082, server_name="0.0.0.0") # share=True generates a public link (requires internet)
    except Exception as e:
        error_logger.error(f"Failed to launch web interface: {e}")
        error_logger.error(f"Traceback: {traceback.format_exc()}")
        raise